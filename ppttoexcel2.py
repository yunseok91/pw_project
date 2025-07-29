from typing import List, Dict
from logger import Logger  # Logger 클래스 import
from PyQt5 import QtWidgets, uic
from PyQt5.QtWidgets import QFileDialog, QMessageBox
from PyQt5.QtCore import QThread, pyqtSignal, QEventLoop, Qt, QTimer
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from PyQt5.QtWidgets import QApplication
from datetime import datetime

import sys
import traceback
import os
import re
import pandas as pd

#pyinstaller -F --noconsole --clean --add-data "pptGuide.ui;." --add-data "logger.py;." --icon="logo.ico" --name "ppttoexcelV2" ppttoexcel2.py


class PPTDataExtractor(QThread):
    # 시그널 정의
    progress_updated = pyqtSignal(int)
    log_message = pyqtSignal(str, str)
    extraction_completed = pyqtSignal(str)
    extraction_error = pyqtSignal(str)

    def __init__(self, ppt_path, logger):
        super().__init__()
        self.ppt_path = ppt_path
        self.logger = logger
        self.excel_output_path = os.path.join(
            os.path.dirname(self.ppt_path),
            f"{os.path.splitext(os.path.basename(self.ppt_path))[0]}_tagging.xlsx"
        )

    def run(self):
        try:
            # 시작 시간 기록
            self.start_time = datetime.now()
            self.log_message.emit(f"시작: {self.start_time.strftime('%Y년 %m월 %d일 %H시 %M분 %S초')}", "info")
            
            # PPT 파일 로드
            presentation = Presentation(self.ppt_path)
            self.log_message.emit(f"PPT 파일 로드 완료. 총 {len(presentation.slides)}개 슬라이드", "info")
            
            # 테이블 데이터 추출
            table_data = self.extract_table_data(presentation)
            
            if not table_data:
                self.log_message.emit("추출된 데이터가 없습니다.", "warning")
                self.extraction_error.emit("추출된 태깅 데이터가 없습니다.")
                return
            
            # 엑셀 파일 경로 설정 (파일명용 시간 - 저장 전에 미리 생성)
            temp_time_str = datetime.now().strftime("%Y%m%d_%H%M%S")
            self.excel_output_path = os.path.join(
                os.path.dirname(self.ppt_path),
                f"{os.path.splitext(os.path.basename(self.ppt_path))[0]}_tagging_{temp_time_str}.xlsx"
            )
            
            # 엑셀 저장
            self.log_message.emit("엑셀 파일 저장 중...", "info")
            self.save_to_excel(table_data)
            
            # 모든 작업 완료 후)
            self.end_time = datetime.now()
            
            # 소요 시간 계산 
            duration = self.end_time - self.start_time
            total_seconds = duration.total_seconds()
            
            # 시간 포맷팅 (소수점 포함)
            if total_seconds < 1:
                time_str = f"{total_seconds:.3f}초"
            elif total_seconds < 60:
                time_str = f"{total_seconds:.2f}초"
            else:
                hours = int(total_seconds // 3600)
                minutes = int((total_seconds % 3600) // 60)
                seconds = total_seconds % 60
                
                if hours > 0:
                    time_str = f"{hours}시간 {minutes}분 {seconds:.2f}초"
                else:
                    time_str = f"{minutes}분 {seconds:.2f}초"
            
            
            # 완료 로그
            self.log_message.emit(f"종료: {self.end_time.strftime('%Y년 %m월 %d일 %H시 %M분 %S초')}", "info")
            self.log_message.emit(f"총 소요시간: {time_str}", "normal")
            
            # 프로그레스 완료 및 결과 전달
            self.progress_updated.emit(100)
            self.extraction_completed.emit(self.excel_output_path)
            
        except Exception as e:
            # 예외 발생 시에도 종료 시간 기록
            if not hasattr(self, 'end_time'):
                self.end_time = datetime.now()
                
            # 부분 완료 시간 계산 (에러 발생 시)
            if hasattr(self, 'start_time'):
                duration = self.end_time - self.start_time
                total_seconds = duration.total_seconds()
                self.log_message.emit(f"오류 발생 전까지 소요시간: {total_seconds:.3f}초", "warning")
            
            self.log_message.emit(f"PPT 변환 오류: {str(e)}", "error")
            self.extraction_error.emit(str(e))

    def extract_table_data(self, presentation) -> List[Dict]:
        table_data = []
        total_slides = len(presentation.slides)
        
        for slide_idx, slide in enumerate(presentation.slides):
            progress = int((slide_idx + 1) / total_slides * 90)
            self.progress_updated.emit(progress)
            
            # 슬라이드 제목 추출
            slide_title = self.extract_slide_title(slide)
            self.log_message.emit(f"슬라이드 {slide_idx + 1}: {slide_title}", "info")
            
            # 테이블 찾기 및 데이터 추출
            for shape in slide.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    table = shape.table
                    
                    # 테이블이 태깅 가이드 형식인지 확인
                    if self.is_tagging_guide_table(table):
                        extracted_data = self.extract_tagging_data(table, slide_idx + 1, slide_title)
                        table_data.extend(extracted_data)
                        self.log_message.emit(f"  - {len(extracted_data)}개 태깅 데이터 추출", "success")
        
        return table_data

    def extract_slide_title(self, slide):
        #슬라이드 제목 추출#
        try:
            # 제목 플레이스홀더 확인
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                return slide.shapes.title.text.strip()
            
            # 상단의 텍스트 박스 찾기
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if hasattr(shape, 'top') and shape.top.pt < 100:  # 상단 100pt 이내
                        text = shape.text.strip()
                        if text:
                            return text
            
            return "제목 없음"
        except:
            return "제목 없음"

    def is_tagging_guide_table(self, table):
        #태깅 가이드 테이블인지 확인#
        try:
            # 첫 번째 행에서 'No' 또는 'Tagging Source' 확인
            if len(table.rows) > 0:
                first_row_text = " ".join(cell.text.strip().lower() for cell in table.rows[0].cells)
                return ('no.' in first_row_text or 'no' in first_row_text) and 'tagging' in first_row_text
            return False
        except:
            return False

    def extract_tagging_data(self, table, slide_num, slide_title):
        #테이블에서 태깅 데이터 추출#
        extracted_data = []
        
        try:
            # 모든 행 데이터 수집
            all_rows = []
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip() if cell.text else "")
                all_rows.append(row_data)
            
            # 헤더 행 건너뛰기
            data_start_idx = 0
            for idx, row in enumerate(all_rows):
                if any('tagging' in cell.lower() for cell in row):
                    data_start_idx = idx + 1
                    break
            
            # 데이터 행 처리
            current_no = None
            current_group_data = []
            
            for row_idx in range(data_start_idx, len(all_rows)):
                row = all_rows[row_idx]
                if not any(row):  # 빈 행
                    continue
                
                # No 확인 (첫 번째 열)
                first_cell = row[0] if row else ""
                
                # 새로운 No 그룹 시작
                if first_cell.isdigit():
                    # 이전 그룹 처리
                    if current_no and current_group_data:
                        extracted_data.extend(self.process_group_data(
                            current_no, current_group_data, slide_num, slide_title
                        ))
                    
                    # 새 그룹 시작
                    current_no = int(first_cell)
                    current_group_data = [row]
                
                # 현재 그룹에 행 추가
                elif current_no is not None:
                    current_group_data.append(row)
            
            # 마지막 그룹 처리
            if current_no and current_group_data:
                extracted_data.extend(self.process_group_data(
                    current_no, current_group_data, slide_num, slide_title
                ))
            
        except Exception as e:
            self.log_message.emit(f"테이블 데이터 추출 오류: {str(e)}", "error")
        
        return extracted_data

    def process_group_data(self, no, group_rows, slide_num, slide_title):
        #No 그룹의 데이터 처리#
        results = []
        
        # Action 찾기
        actions = []
        for row in group_rows:
            for cell in row[1:]:  # 첫 번째 열(No) 제외
                if cell and not any(tag in cell for tag in ['AA', 'GA', 'data-omni', 'ga-']):
                    # Action 키워드 체크
                    if any(keyword in cell.lower() for keyword in 
                          ['click', 'buy', 'order', 'reserve', 'open', 'close', 'drop', 'where', 'pre-order']) or \
                       (len(cell) > 5 and '=' not in cell and cell not in ['　', '']):
                        actions.append(cell)
        
        # 각 Action에 대한 태깅 정보 찾기
        if not actions:  # Action이 없으면 전체를 하나의 항목으로 처리
            actions = ['']
        
        for action in actions:
            result = {
                'No': no,
                'Slide': slide_num,
                'Title': slide_title,
                'Action': action
            }
            
            # 모든 행에서 태깅 정보 추출
            all_text = " ".join(" ".join(row) for row in group_rows)
            
            # 태깅 속성 추출
            tagging_attrs = self.extract_tagging_attributes(all_text)
            result.update(tagging_attrs)
            
            # Action별로 개별 태깅이 있는 경우 처리
            if action:
                for row in group_rows:
                    row_text = " ".join(row)
                    if action in row_text:
                        # 해당 Action이 있는 행 주변의 태깅 정보 추출
                        specific_attrs = self.extract_tagging_attributes(row_text)
                        if specific_attrs:
                            result.update(specific_attrs)
            
            results.append(result)
        
        return results

    def extract_tagging_attributes(self, text):
        #텍스트에서 태깅 속성 추출#
        attributes = {}
        
        # 텍스트 정리
        text = re.sub(r'""', '"', text)  # 이중 따옴표 처리
        text = re.sub(r'“', '"', text)   # 특수 따옴표 처리
        text = re.sub(r'”', '"', text)
        text = re.sub(r'"', '"', text)   # 추가 특수 따옴표 처리

        # 패턴 매칭
        patterns = {
            'data-omni-type': r'data-omni-type\s*=\s*"([^"]+)"',
            'data-omni': r'data-omni\s*=\s*"([^"]+)"',
            'ga-ca': r'ga-ca\s*=\s*"([^"]+)"',
            'ga-ac': r'ga-ac\s*=\s*"([^"]+)"',
            'ga-la': r'ga-la\s*=\s*"([^"]+)"'
        }
        
        for key, pattern in patterns.items():
            match = re.search(pattern, text)
            if match:
                value = match.group(1).strip()
                # 잘못된 값 정리
                if key == 'data-omni-type':
                    value = value.replace('" data-omni=', '')
                attributes[key] = value
        
        return attributes

    def save_to_excel(self, table_data):
        #데이터를 엑셀로 저장#
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = '태깅 데이터'
            
            # 헤더 설정
            headers = ['No', 'Slide', 'Title', 'Action', 'data-omni-type', 'data-omni', 'ga-ca', 'ga-ac', 'ga-la']
            
            # 헤더 스타일
            header_font = Font(bold=True)
            header_fill = PatternFill(start_color="DDEBF7", end_color="DDEBF7", fill_type="solid")
            header_alignment = Alignment(horizontal="center", vertical="center")
            
            # 헤더 작성 (5행, D열부터)
            start_row = 5
            start_col = 4  # D열
            
            for col_idx, header in enumerate(headers):
                cell = ws.cell(row=start_row, column=start_col + col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
            
            # 데이터 작성
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            for row_idx, data in enumerate(table_data, start=start_row + 1):
                for col_idx, header in enumerate(headers):
                    cell = ws.cell(
                        row=row_idx,
                        column=start_col + col_idx,
                        value=data.get(header, '')
                    )
                    cell.border = thin_border
            
            # 파일 저장
            wb.save(self.excel_output_path)
            self.log_message.emit(f"엑셀 파일 저장 완료: {self.excel_output_path}", "success")
            
        except Exception as e:
            self.log_message.emit(f"엑셀 저장 오류: {str(e)}", "error")
            raise

def resource_path(relative_path):
    #PyInstaller 리소스 경로 처리#
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

class PPTConverterApp(QtWidgets.QMainWindow):
    def __init__(self):
        super().__init__()
        
        # UI 파일 로드
        uic.loadUi(resource_path('pptGuide2.ui'), self)        
        # 초기 설정
        self.ppt_file_path = ""
        
        # Logger 초기화
        self.logger = Logger(self.logTextEdit)
        self.logger.log_signal.connect(self.logTextEdit.insertHtml)
        
        # 초기 로그 메시지
        self.logger.log("Unpack Tagging Guide.", "info")
        self.logger.log("PPT 파일을 선택해주세요.", "normal")
        
        # 버튼 이벤트 연결
        self.pptSelectBtn.clicked.connect(self.select_ppt_file)
        self.convertBtn.clicked.connect(self.convert_ppt)
        self.clearLogBtn.clicked.connect(self.clear_log)
        self.exitBtn.clicked.connect(self.close)
        
        # 초기에는 변환 버튼 비활성화
        self.convertBtn.setEnabled(False)
        
        # 타이머 초기화
        self.update_timer = None

    def select_ppt_file(self):
        #PPT 파일 선택#
        file_path, _ = QFileDialog.getOpenFileName(
            self, 
            "PowerPoint 파일 선택", 
            "", 
            "PowerPoint 파일 (*.pptx *.ppt)"
        )
        
        if file_path:
            self.ppt_file_path = file_path
            self.pptPathLabel.setText(os.path.basename(file_path))
            self.logger.log(f"PPT 파일이 선택되었습니다: {os.path.basename(file_path)}", "success")
            self.convertBtn.setEnabled(True)
    
    def convert_ppt(self):
        #PPT 변환#
        if not self.ppt_file_path:
            QMessageBox.warning(self, "경고", "PPT 파일을 먼저 선택해주세요.")
            return
        
        self.logger.log("PPT 변환을 시작합니다...", "normal")
        self.convertBtn.setEnabled(False)
        self.progressBar.setValue(0)
        
        # 데이터 추출기 생성
        self.extractor = PPTDataExtractor(self.ppt_file_path, self.logger)
        
        # 시그널 연결
        self.extractor.progress_updated.connect(self.update_progress)
        self.extractor.log_message.connect(self.handle_log_message)
        self.extractor.extraction_completed.connect(self.conversion_finished)
        self.extractor.extraction_error.connect(self.conversion_error)
        
        # 타이머 설정
        if self.update_timer:
            self.update_timer.stop()
        self.update_timer = QTimer(self)
        self.update_timer.timeout.connect(self.process_events)
        self.update_timer.start(100)
        
        # 추출 스레드 시작
        self.extractor.start()
    
    def handle_log_message(self, message, msg_type):
        #로그 메시지 처리#
        self.logger.log(message, msg_type)
        self.scroll_log_to_bottom()
    
    def update_progress(self, value):
        #진행률 업데이트#
        self.progressBar.setValue(value)
        self.scroll_log_to_bottom()
    
    def scroll_log_to_bottom(self):
        #로그 스크롤을 맨 아래로#
        scrollbar = self.logTextEdit.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())
    
    def process_events(self):
        #주기적으로 이벤트 처리#
        QApplication.processEvents(QEventLoop.ExcludeUserInputEvents)
        self.scroll_log_to_bottom()
    
    def conversion_finished(self, excel_path):
        #변환 완료 처리#
        self.convertBtn.setEnabled(True)
        
        if self.update_timer:
            self.update_timer.stop()
        
        self.logger.log("PPT 변환 완료!", "normal")
        
        reply = QMessageBox.question(
            self, 
            "변환 완료",
            f"Excel 파일이 생성되었습니다.\n{excel_path}\n\n파일을 열어보시겠습니까?",
            QMessageBox.Yes | QMessageBox.No
        )
        
        if reply == QMessageBox.Yes:
            try:
                if sys.platform == 'win32':
                    os.startfile(excel_path)
                elif sys.platform == 'darwin':
                    os.system(f'open "{excel_path}"')
                else:
                    os.system(f'xdg-open "{excel_path}"')
            except Exception as e:
                QMessageBox.warning(self, "오류", f"파일을 열 수 없습니다: {str(e)}")
    
    def conversion_error(self, error_msg):
        #변환 오류 처리#
        self.convertBtn.setEnabled(True)
        
        if self.update_timer:
            self.update_timer.stop()
        
        QMessageBox.critical(self, "변환 오류", error_msg)
    
    def clear_log(self):
        #로그 지우기#
        self.logTextEdit.clear()
        self.logger.log("로그가 지워졌습니다.", "info")


def main():
    try:
        app = QtWidgets.QApplication(sys.argv)
        window = PPTConverterApp()
        window.show()
        sys.exit(app.exec_())
        
    except Exception as e:
        # 충돌 시 로그 파일 생성
        crash_log_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "crash_logs")
        if not os.path.exists(crash_log_dir):
            os.makedirs(crash_log_dir)
        
        crash_log_path = os.path.join(
            crash_log_dir,
            f"main_crash_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        )
        
        with open(crash_log_path, 'w', encoding='utf-8') as f:
            f.write(f"프로그램 충돌: {str(e)}\n\n")
            f.write(f"Timestamp: {datetime.now()}\n")
            f.write("Stack Trace:\n")
            traceback.print_exc(file=f)
        
        print(f"프로그램이 충돌했습니다. 로그가 저장되었습니다: {crash_log_path}")


if __name__ == "__main__":
    main()