from typing import List, Dict
from logger import Logger  # Logger 클래스 import
from PyQt5 import QtWidgets, uic
from PyQt5.QtWidgets import QFileDialog, QMessageBox
from PyQt5.QtCore import QThread, pyqtSignal
from PyQt5.QtCore import QThread, pyqtSignal, QEventLoop, Qt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.text.text import _Run
from PyQt5.QtWidgets import QApplication

import sys
import datetime 
import traceback
import os
import re
import pandas as pd

class PPTDataExtractor(QThread):
    #Qthread 상속 run
    #시그널 정의 pyqtSignal 통신채널 
    progress_updated = pyqtSignal(int) #진행률 전달
    log_message = pyqtSignal(str, str)  # 메시지, 타입
    extraction_completed = pyqtSignal(str) #완료 시 파일 경로 전달
    extraction_error = pyqtSignal(str) #오류 메시지 전달

    def __init__(self,ppt_path,logger):
        super().__init__()
        self.ppt_path = ppt_path
        self.logger = logger  # 전달받은 logger 그대로 사용
        #출력 파일 경로 
        self.excel_output_path = os.path.join(
            os.path.dirname(self.ppt_path),
            f"{os.path.splitext(os.path.basename(self.ppt_path))[0]}_tagging.xlsx"
        )

    def run(self):
        try:
            #ppt 파일 로드 
            presentation = Presentation(self.ppt_path)
            self.log_message.emit(f"PPT 파일 로드 완료. 총 {len(presentation.slides)}개 슬라이드", "info")
            
            #테이블 데이터
            table_data = self.extract_table_data(presentation)

            #엑셀 저장
            self.save_to_excel(table_data)
            
            self.progress_updated.emit(100)
            #완료시그널
            self.extraction_completed.emit(self.excel_output_path)
            #작업중 시그널 전송 emit 데이터 전송
        except Exception as e : 
            self.log_message.emit(f"PPT변환 오류 {e}","error")
            self.extraction_error.emit(str(e))
            
    def extract_table_data(self, presentation) -> List[Dict]:
        table_data = []
        try:
            # Python의 가비지 컬렉션(Garbage Collection)
            import gc

            self.log_message.emit("테이블 데이터 추출 시작", "info")
            total_slides = len(presentation.slides)
            
            # 메모리 사용량 모니터링 변수 추가
            processed_shapes = 0
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                progress = int((slide_idx / total_slides) * 100)
                self.progress_updated.emit(progress)
                QApplication.processEvents(QEventLoop.ExcludeUserInputEvents)

                try:
                    slide_title = f"슬라이드 {slide_idx}"  # 기본값
                    top_threshold = 100.0
                    tables_found = False

                    # 한 번의 루프로 슬라이드 제목과 테이블 모두 처리
                    for shape in slide.shapes:
                        processed_shapes += 1
                        
                        # 슬라이드 제목 찾기 (아직 발견하지 못한 경우만)
                        if slide_title == f"슬라이드 {slide_idx}" and hasattr(shape, "text_frame") and shape.text_frame:
                            top = shape.top.pt
                            if top <= top_threshold:  # 상단에 있는 텍스트 도형만 제목으로 간주
                                text = shape.text.strip()
                                if text:
                                    slide_title = text
                        
                        # 테이블 처리
                        if hasattr(shape, 'has_table') and shape.has_table:
                            tables_found = True
                            table = shape.table
                            
                            # 빈 테이블 확인
                            if len(table.rows) <= 1:
                                self.log_message.emit(f"슬라이드 {slide_idx}에 빈 테이블이 있습니다.", "warning")
                                continue
                                
                            # 헤더 확인
                            first_row_cells = table.rows[0].cells
                            if len(first_row_cells) == 0:
                                self.log_message.emit(f"슬라이드 {slide_idx}의 테이블에 헤더가 없습니다.", "warning")
                                continue
                                
                            headers = [cell.text.strip() for cell in first_row_cells]

                            # 필요한 헤더가 있는지 확인
                            if 'No.' in headers and 'Tagging Source' in headers:
                                header_index_no = headers.index('No.')
                                header_index_action = header_index_no + 1  # 'Action' 컬럼 위치 추정
                                
                                for row_idx in range(1, len(table.rows)):
                                    row_cells = table.rows[row_idx].cells
                                    
                                    # 행에 셀이 충분히 있는지 확인
                                    if len(row_cells) <= max(header_index_no, header_index_action):
                                        self.log_message.emit(f"슬라이드 {slide_idx}, 행 {row_idx}에 셀이 부족합니다.", "warning")
                                        continue
                                    
                                    try:
                                        no = row_cells[header_index_no].text.strip()
                                        
                                        # Action 값 안전하게 가져오기
                                        try:
                                            if len(row_cells) > 1:
                                                full_text = row_cells[1].text.strip()
                                                # 줄바꿈이 있는 경우 첫 번째 줄만 사용
                                                action = full_text.split("\n")[0] if "\n" in full_text else full_text
                                            else:
                                                action = ""
                                        except Exception as action_error:
                                            self.log_message.emit(f"슬라이드 {slide_idx}, 행 {row_idx}의 Action 추출 중 오류: {action_error}", "warning")
                                            action = ""
                                        
                                        # 모든 셀의 텍스트를 결합 (오류에 강하게)
                                        row_text = ' '.join(cell.text for cell in row_cells if hasattr(cell, 'text'))

                                        patterns = {
                                            'data-omni-type': r'data-omni-type="([^"]*)"',
                                            'data-omni': r'data-omni="([^"]*)"',
                                            'ga-ca': r'ga-ca="([^"]*)"',
                                            'ga-ac': r'ga-ac="([^"]*)"',
                                            'ga-la': r'ga-la="([^"]*)"'
                                        }

                                        extracted_data = {}
                                        for key, pattern in patterns.items():
                                            match = re.search(pattern, row_text)
                                            if match:
                                                extracted_data[key] = match.group(1)

                                        table_entry = {
                                            'Slide': slide_idx,
                                            'Title': slide_title,
                                            'No': no,
                                            'Action': action,
                                            **extracted_data
                                        }
                                        table_data.append(table_entry)
                                    except Exception as cell_error:
                                        self.log_message.emit(f"슬라이드 {slide_idx}, 행 {row_idx} 처리 중 셀 오류: {cell_error}", "warning")
                                        continue
                    
                    if not tables_found:
                        self.log_message.emit(f"슬라이드 {slide_idx}에 테이블이 없습니다.", "info")

                except Exception as slide_error:
                    self.log_message.emit(f"슬라이드 {slide_idx} 처리 중 오류: {slide_error}", "warning")
                    continue

                # 메모리 관리 개선: 더 자주 GC 호출 (20개 슬라이드마다)
                if slide_idx % 20 == 0 or processed_shapes > 500:
                    gc.collect()
                    processed_shapes = 0
                    self.log_message.emit(f"메모리 정리 수행 (슬라이드 {slide_idx}/{total_slides})", "info")
                    QApplication.processEvents(QEventLoop.ExcludeUserInputEvents)

            # 데이터 추출 결과 확인
            if not table_data:
                self.log_message.emit("추출된 테이블 데이터가 없습니다. PPT 형식을 확인해주세요.", "warning")
            else:
                self.log_message.emit(f"총 추출된 데이터 수: {len(table_data)}", "success")

        except Exception as e:
            self.log_message.emit(f"테이블 데이터 추출 중 심각한 오류: {e}", "error")
            import traceback
            self.log_message.emit(traceback.format_exc(), "error")

        return table_data

    
    def save_to_excel(self, table_data):
        try:
            # 데이터 로깅
            self.log_message.emit(f"총 추출된 데이터 개수: {len(table_data)}", "info")
            
            # 데이터 내용 상세 로깅
            for idx, row in enumerate(table_data, 1):
                self.log_message.emit(f"행 {idx} 데이터: {row}", "warning")

            # 데이터프레임 생성
            df_table = pd.DataFrame(table_data)

            # 컬럼 순서 지정
            columns_order = [
                'Slide', 'No','Title', 'Action', 
                'data-omni-type', 'data-omni', 
                'ga-ca', 'ga-ac', 'ga-la'
            ]

            # 존재하는 컬럼만 선택
            columns_order = [col for col in columns_order if col in df_table.columns]
            
            # 데이터프레임 재정렬
            df_table = df_table[columns_order]

            # 워크북과 워크시트 생성
            wb = Workbook()
            ws = wb.active
            ws.title = '테이블 데이터'

            # 헤더 스타일 설정
            header_font = Font(bold=True)
            header_fill = PatternFill(start_color="DDEBF7", end_color="DDEBF7", fill_type="solid")
            header_alignment = Alignment(horizontal="center", vertical="center")
            header_border = Border(
                left=Side(style='thin'), 
                right=Side(style='thin'), 
                top=Side(style='thin'), 
                bottom=Side(style='thin')
            )

            # 헤더 작성 (2행부터)
            for col, header in enumerate(columns_order, start=1):
                cell = ws.cell(row=2, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
                cell.border = header_border

            # 데이터 작성 (3행부터)
            for row_idx, row_data in enumerate(df_table.to_dict('records'), start=3):
                for col_idx, header in enumerate(columns_order, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=row_data.get(header, ''))

            # 고유한 파일 이름 생성
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            self.excel_output_path = os.path.join(
                os.path.dirname(self.ppt_path),
                f"{os.path.splitext(os.path.basename(self.ppt_path))[0]}_tagging_{timestamp}.xlsx"
            )

            # 엑셀 파일 저장
            wb.save(self.excel_output_path)

            self.log_message.emit(f"엑셀 파일 저장 완료: {self.excel_output_path}", "success")

        except Exception as e:
            self.log_message.emit(f"엑셀 저장 중 오류: {e}", "error")
            import traceback
            self.log_message.emit(traceback.format_exc(), "error")
            raise


class PPTConverterApp(QtWidgets.QMainWindow):
    def __init__(self):
        # 부모 클래스 초기화
        super().__init__()
        
        # UI 파일 로드
        uic.loadUi('pptGuide.ui', self)
        
        # 초기 설정
        self.ppt_file_path = ""
        # Logger 초기화 
        
        self.logger = Logger(self.logTextEdit)
        self.logger.log_signal.connect(self.logTextEdit.insertHtml)

        # 초기 로그 메시지
        self.logger.log("Unpack Tagging Guide.", "info")
        self.logger.log("PPT 파일을 선택해주세요.", "normal")
        # 버튼 이벤트 연결
        self.pptSelectBtn.clicked.connect(self.select_ppt_file)
        self.convertBtn.clicked.connect(self.convert_ppt)
        self.clearLogBtn.clicked.connect(self.clear_log)
        self.exitBtn.clicked.connect(self.close)
        
        # 초기에는 변환 버튼 비활성화
        self.convertBtn.setEnabled(False)

    def select_ppt_file(self):
        #PPT 파일 선택 함수#
        # 파일 대화상자 열기
        file_path, _ = QFileDialog.getOpenFileName(
            self, 
            "PowerPoint 파일 선택", 
            "", 
            "PowerPoint 파일 (*.pptx *.ppt)"
        )
        
        # 파일 선택되었으면
        if file_path:
            # 파일 경로 저장
            self.ppt_file_path = file_path
            
            # 선택된 파일 경로 표시
            self.pptPathLabel.setText(os.path.basename(file_path))
            self.logger.log(f"PPT 파일이 선택되었습니다: {os.path.basename(file_path)}", "success")
            # 변환 버튼 활성화
            self.convertBtn.setEnabled(True)
    
    def convert_ppt(self):
        #PPT 변환 함수#
        # 파일 경로 확인
        if not self.ppt_file_path:
            QMessageBox.warning(self, "경고", "PPT 파일을 먼저 선택해주세요.")
            return
        
        # 변환 시작 로그 표시
        self.logger.log("PPT 변환을 시작합니다...","normal")
        # 변환 버튼 비활성화
        self.convertBtn.setEnabled(False)

        # 데이터 추출기 생성
        self.extractor = PPTDataExtractor(self.ppt_file_path, self.logger)
        # 시그널 연결
        self.extractor.progress_updated.connect(self.update_progress)
        self.extractor.log_message.connect(self.handle_log_message)  # 로그 메시지 처리 함수 연결
        self.extractor.extraction_completed.connect(self.conversion_finished)
        self.extractor.extraction_error.connect(self.conversion_error)
            # 주기적인 이벤트 루프 처리를 위한 타이머 설정
        from PyQt5.QtCore import QTimer
        self.update_timer = QTimer(self)
        self.update_timer.timeout.connect(self.process_events)
        self.update_timer.start(100)  # 100ms마다 이벤트 처리
    
        # 추출 스레드 시작
        self.extractor.start()
    def handle_log_message(self, message, msg_type):
    # 로그 메시지를 받아서 로거로 전달하고 UI 업데이트
        self.logger.log(message, msg_type)
        self.scroll_log_to_bottom()
    def update_progress(self, value):
        self.progressBar.setValue(value)
        self.scroll_log_to_bottom()
    def scroll_log_to_bottom(self):
     # 로그 텍스트를 항상 맨 아래로 스크롤
        scrollbar = self.logTextEdit.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())  
    def process_events(self):
        # 주기적으로 이벤트 루프 처리
        QApplication.processEvents(QEventLoop.ExcludeUserInputEvents)
        self.scroll_log_to_bottom()

    def conversion_finished(self,excel_path):
        #변환 활성화
        self.convertBtn.setEnabled(True)

        #완료 
        self.logger.log("PPT변환 완료!","success")
        #엑셀 열기
        reply = QMessageBox.question(
            self,"변환 완료",
            f"Excel 파일이 생성되었습니다.\n{excel_path}\n 파일을 열어보시겠습니까?",
            QMessageBox.Yes | QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            try:
                os.startfile(excel_path)
            except Exception as e:
                QMessageBox.warning(self, "오류", f"파일을 열 수 없습니다: {str(e)}")
    def conversion_error(self, error_msg):
    # 변환 버튼 다시 활성화
        self.convertBtn.setEnabled(True)
    
        # 오류 메시지 표시
        QMessageBox.critical(self, "변환 오류", error_msg)
   
    def clear_log(self):
        #로그 지우기 함수#
        self.logTextEdit.clear()
        self.logger.log("로그가 지워졌습니다.", "info")
    
def main():
    try :
        # Qt 애플리케이션 생성
        app = QtWidgets.QApplication(sys.argv)
        
        # 메인 윈도우 생성
        window = PPTConverterApp()
        window.show()
        
        # 애플리케이션 실행
        sys.exit(app.exec_())

    except Exception as e : 
        # 충돌 시 로그 파일 생성
        crash_log_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)),"crash_log")
        if not os.path.exists(crash_log_dir):
            os.makedirs(crash_log_dir)

        crash_log_path = os.path.join(
            crash_log_dir,
            f"main_crash_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        )
        #충돌 로그를 쓸 때 파일을 열고 기록한 후 자동으로 닫기
        with open(crash_log_path, 'w', encoding='utf-8') as f :
            f.write(f"프로그램 충돌:{str(e)}\n\n")
            f.write(f"Timestamp {datetime.datetime.now()}\n")
            f.write("Stavck Trace:\n")
            traceback.print_exc(file=f)
    print(f"프로그램이 충돌했습니다. 로그가 저장되었습니다: {crash_log_path}")
        
if __name__ == "__main__":
    main()
